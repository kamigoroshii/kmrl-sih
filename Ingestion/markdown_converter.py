import fitz
import docx
import os
import re
from datetime import datetime
import tempfile
import subprocess
import sys
import base64
from PIL import Image
import io
import easyocr
from Ingestion.clip_embedder import embed_image_clip

# Initialize EasyOCR reader
reader = easyocr.Reader(['en'])

def extract_image_content(img_path):
    """
    Extract OCR text and CLIP embedding from an image.
    Returns a tuple of (ocr_text, clip_embedding).
    """
    ocr_text = ""
    clip_embedding = None
    
    try:
        # Perform OCR on the image
        ocr_result = reader.readtext(img_path, detail=0)
        if isinstance(ocr_result, list):
            ocr_text = '\n'.join([str(r) for r in ocr_result])
        else:
            ocr_text = str(ocr_result)
        
        print(f"   📝 OCR extracted {len(ocr_text)} characters from image")
        
    except Exception as e:
        print(f"   ⚠️  OCR failed for image: {e}")
        ocr_text = ""
    
    try:
        # Generate CLIP embedding
        clip_embedding = embed_image_clip(img_path)
        print(f"   🎯 CLIP embedding generated ({len(clip_embedding)} dimensions)")
        
    except Exception as e:
        print(f"   ⚠️  CLIP embedding failed for image: {e}")
        clip_embedding = None
    
    return ocr_text, clip_embedding

def pdf_to_markdown(filepath):
    """
    Convert PDF to Markdown format.
    Handles both digital PDFs (with text) and scanned PDFs (using OCR).
    Also extracts and handles images within the PDF with OCR and CLIP embeddings.
    """
    doc = fitz.open(filepath)
    markdown_content = []
    
    # Extract document metadata
    metadata = doc.metadata
    title = metadata.get('title', '') if metadata else ''
    author = metadata.get('author', '') if metadata else ''
    
    # Add document header
    if title:
        markdown_content.append(f"# {title}\n")
    if author:
        markdown_content.append(f"**Author:** {author}\n")
    
    markdown_content.append(f"**Source:** {os.path.basename(filepath)}\n")
    markdown_content.append(f"**Converted:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
    
    # Create images directory for extracted images
    base_dir = os.path.dirname(filepath)
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    images_dir = os.path.join(base_dir, f"{base_name}_images")
    os.makedirs(images_dir, exist_ok=True)
    
    for page_num, page in enumerate(doc, 1):
        markdown_content.append(f"## Page {page_num}\n")
        
        # Extract images from the page
        image_list = page.get_images()
        image_count = 0
        
        for img_index, img in enumerate(image_list):
            try:
                # Get image data
                xref = img[0]
                pix = fitz.Pixmap(doc, xref)
                
                # Handle different color spaces
                if pix.n - pix.alpha < 4:  # GRAY or RGB
                    img_data = pix.tobytes("png")
                else:  # CMYK or other color spaces: convert to RGB first
                    try:
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix1.tobytes("png")
                        pix1 = None
                    except Exception as cmyk_error:
                        # If RGB conversion fails, try to save as JPEG instead
                        try:
                            img_data = pix.tobytes("jpeg")
                        except Exception as jpeg_error:
                            print(f"[WARNING] Failed to convert image {img_index + 1} from page {page_num}: {cmyk_error}")
                            pix = None
                            continue
                
                # Save image to file
                img_filename = f"page_{page_num}_img_{img_index + 1}.png"
                img_path = os.path.join(images_dir, img_filename)
                
                with open(img_path, "wb") as img_file:
                    img_file.write(img_data)
                
                # Extract OCR text and CLIP embedding from the image
                print(f"   🔍 Processing image {img_index + 1} from page {page_num}...")
                ocr_text, clip_embedding = extract_image_content(img_path)
                
                # Add image reference to markdown
                markdown_content.append(f"![Image {img_index + 1}]({img_path})\n\n")
                
                # Add OCR text if available
                if ocr_text.strip():
                    markdown_content.append(f"**OCR Text from Image {img_index + 1}:**\n")
                    markdown_content.append(f"```\n{ocr_text.strip()}\n```\n\n")
                
                # Add CLIP embedding metadata if available
                if clip_embedding:
                    markdown_content.append(f"**CLIP Embedding (Image {img_index + 1}):**\n")
                    markdown_content.append(f"*[CLIP visual embedding generated - {len(clip_embedding)} dimensions]*\n\n")
                
                image_count += 1
                
                pix = None
                
            except Exception as e:
                print(f"[WARNING] Failed to extract image {img_index + 1} from page {page_num}: {e}")
                continue
        
        # Try to extract text (for digital PDFs)
        text = page.get_text().strip()
        
        if not text:
            # If no text found, it might be a scanned PDF
            if image_count == 0:
                markdown_content.append("*[Scanned page - text extraction not available]*\n\n")
            else:
                markdown_content.append("*[Page contains images only]*\n\n")
            continue
        
        # Process the text to create markdown structure
        lines = text.split('\n')
        processed_lines = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
                
            # Detect headers (lines that are shorter and might be titles)
            if len(line) < 100 and len(line) > 3:
                # Check if it looks like a header (all caps, ends with numbers, etc.)
                if (line.isupper() or 
                    re.match(r'^[A-Z][A-Za-z\s]+\d*$', line) or
                    re.match(r'^\d+\.\s+[A-Z]', line)):
                    processed_lines.append(f"### {line}\n")
                    continue
            
            # Detect lists
            if re.match(r'^\d+\.\s+', line):
                processed_lines.append(line + "\n")
                continue
            elif re.match(r'^[•\-\*]\s+', line):
                processed_lines.append(line + "\n")
                continue
            
            # Regular paragraph
            processed_lines.append(line + "\n")
        
        if processed_lines:
            markdown_content.extend(processed_lines)
        
        markdown_content.append("\n")
    
    doc.close()
    return "\n".join(markdown_content)

def docx_to_markdown(filepath):
    """
    Convert DOCX to Markdown format.
    """
    doc = docx.Document(filepath)
    markdown_content = []
    
    # Extract document properties
    core_props = doc.core_properties
    title = core_props.title if core_props and core_props.title else ''
    author = core_props.author if core_props and core_props.author else ''
    
    # Add document header
    if title:
        markdown_content.append(f"# {title}\n")
    if author:
        markdown_content.append(f"**Author:** {author}\n")
    
    markdown_content.append(f"**Source:** {os.path.basename(filepath)}\n")
    markdown_content.append(f"**Converted:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
    
    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        
        # Handle different paragraph styles
        style = paragraph.style.name.lower() if paragraph.style and paragraph.style.name else ''
        
        if 'heading' in style or 'title' in style:
            # Determine heading level based on style
            if '1' in style or 'title' in style:
                markdown_content.append(f"# {text}\n")
            elif '2' in style:
                markdown_content.append(f"## {text}\n")
            elif '3' in style:
                markdown_content.append(f"### {text}\n")
            else:
                markdown_content.append(f"## {text}\n")
        else:
            # Regular paragraph
            markdown_content.append(f"{text}\n\n")
    
    # Handle tables
    for table in doc.tables:
        markdown_content.append("\n")
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip() for cell in row.cells]
            if i == 0:
                # Header row
                markdown_content.append("| " + " | ".join(cells) + " |\n")
                markdown_content.append("| " + " | ".join(["---"] * len(cells)) + " |\n")
            else:
                # Data row
                markdown_content.append("| " + " | ".join(cells) + " |\n")
        markdown_content.append("\n")
    
    return "\n".join(markdown_content)

def pptx_to_markdown(filepath):
    """
    Convert PPTX to Markdown format.
    Note: This requires python-pptx library which may not be installed.
    """
    try:
        from pptx import Presentation
    except ImportError:
        # Fallback: try to install python-pptx or use alternative method
        print(f"[WARNING] python-pptx not available. Cannot convert {filepath} to markdown.")
        return f"# {os.path.basename(filepath)}\n\n*[PowerPoint file - conversion not available]*\n\n"
    
    try:
        prs = Presentation(filepath)
        markdown_content = []
        
        # Extract presentation properties
        core_props = prs.core_properties
        title = core_props.title if core_props and core_props.title else ''
        author = core_props.author if core_props and core_props.author else ''
        
        # Add presentation header
        if title:
            markdown_content.append(f"# {title}\n")
        if author:
            markdown_content.append(f"**Author:** {author}\n")
        
        markdown_content.append(f"**Source:** {os.path.basename(filepath)}\n")
        markdown_content.append(f"**Converted:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
        
        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_content.append(f"## Slide {slide_num}\n")
            
            for shape in slide.shapes:
                # Check if shape has text attribute and it's not None
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Try to detect if this is a title
                    shape_name = getattr(shape, 'name', '')
                    if shape_name and 'title' in shape_name.lower():
                        markdown_content.append(f"### {text}\n")
                    else:
                        markdown_content.append(f"{text}\n\n")
            
            markdown_content.append("\n")
        
        return "\n".join(markdown_content)
    
    except Exception as e:
        print(f"[ERROR] Failed to convert PPTX {filepath}: {e}")
        return f"# {os.path.basename(filepath)}\n\n*[PowerPoint file - conversion failed]*\n\n"

def convert_to_markdown(filepath):
    """
    Main function to convert various file types to Markdown.
    Returns the markdown content as a string.
    """
    ext = os.path.splitext(filepath)[1].lower()
    
    try:
        if ext == '.pdf':
            return pdf_to_markdown(filepath)
        elif ext == '.docx':
            return docx_to_markdown(filepath)
        elif ext == '.pptx':
            return pptx_to_markdown(filepath)
        else:
            print(f"[WARNING] Unsupported file type for markdown conversion: {ext}")
            return None
    except Exception as e:
        print(f"[ERROR] Failed to convert {filepath} to markdown: {e}")
        return None

def save_as_markdown(filepath, output_dir=None):
    """
    Convert a file to markdown and save it to a .md file.
    Returns the path to the saved markdown file.
    """
    if output_dir is None:
        output_dir = os.path.dirname(filepath)
    
    markdown_content = convert_to_markdown(filepath)
    if markdown_content is None:
        return None
    
    # Create markdown filename
    base_name = os.path.splitext(os.path.basename(filepath))[0]
    markdown_filename = f"{base_name}.md"
    markdown_path = os.path.join(output_dir, markdown_filename)
    
    try:
        with open(markdown_path, 'w', encoding='utf-8') as f:
            f.write(markdown_content)
        print(f"✅ Converted {filepath} to {markdown_path}")
        return markdown_path
    except Exception as e:
        print(f"[ERROR] Failed to save markdown file {markdown_path}: {e}")
        return None 